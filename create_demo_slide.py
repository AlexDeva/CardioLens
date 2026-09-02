from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT = "CardioLens_Demo_Slide.pptx"
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Palette
navy = RGBColor(10, 22, 38)
ink = RGBColor(22, 35, 52)
muted = RGBColor(166, 184, 199)
white = RGBColor(248, 251, 253)
coral = RGBColor(255, 103, 96)
cyan = RGBColor(48, 207, 190)
yellow = RGBColor(247, 193, 76)
blue = RGBColor(89, 151, 255)
violet = RGBColor(157, 125, 255)
panel = RGBColor(19, 38, 59)
panel_2 = RGBColor(27, 49, 70)
line = RGBColor(48, 71, 92)

# Background
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = navy


def text(x, y, w, h, value, size=14, color=white, bold=False, align=PP_ALIGN.LEFT, font="Aptos"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = value
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def rect(x, y, w, h, fill, radius=False, border=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = border or fill
    if radius:
        s.adjustments[0] = 0.12
    return s


def line_shape(x1, y1, x2, y2, color, width=1.5, arrow=False):
    s = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    s.line.color.rgb = color
    s.line.width = Pt(width)
    if arrow:
        s.line.end_arrowhead = True
    return s


connector = line_shape


def pill(x, y, w, label, fill, fg=navy):
    rect(x, y, w, 0.28, fill, radius=True)
    text(x, y + 0.005, w, 0.26, label, 9, fg, True, PP_ALIGN.CENTER)


def add_circle(x, y, diameter, fill, border=None, width=1.2):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(diameter), Inches(diameter)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border or fill
    shape.line.width = Pt(width)
    return shape

# Header
text(0.65, 0.32, 8.2, 0.42, "CardioLens | DEMO IMPLEMENTATION", 13, cyan, True)
text(0.65, 0.78, 11.9, 0.62, "Capture the profile. Trace the risk. Act with evidence.", 25, white, True)
text(0.67, 1.42, 11.5, 0.3, "A live Streamlit decision-support workflow for research and education", 11, muted)

# Left: pseudo-3D pipeline
text(0.7, 1.95, 5.2, 0.28, "THE CASCADE ENGINE", 10, yellow, True)
text(0.7, 2.2, 5.4, 0.32, "Three connected model layers", 17, white, True)

# 3D blocks: shadow/depth layers and front face
blocks = [
    (0.78, 2.82, 3.72, 0.9, coral, "01  METABOLIC LAYER", "XGBoost • obesity probability"),
    (1.26, 3.93, 3.72, 0.9, cyan, "02  CARDIOVASCULAR LAYER", "LightGBM • 10-year CVD risk"),
    (1.74, 5.04, 3.72, 0.9, yellow, "03  EXPLAINABILITY LAYER", "Feature drivers • action plan"),
]
for x, y, w, h, color, heading, sub in blocks:
    # extrusion gives the blocks a 3D/isometric feel
    rect(x + 0.13, y + 0.13, w, h, RGBColor(5, 13, 24), radius=True)
    rect(x + 0.06, y + 0.06, w, h, color, radius=True)
    rect(x, y, w, h, panel, radius=True, border=color)
    rect(x, y, 0.11, h, color, radius=True)
    text(x + 0.25, y + 0.13, w - 0.38, 0.25, heading, 11, color, True)
    text(x + 0.25, y + 0.45, w - 0.38, 0.24, sub, 10, white)

# Connecting rails
line_shape(2.64, 3.72, 2.64, 3.91, muted, 1.5, True)
line_shape(3.12, 4.83, 3.12, 5.02, muted, 1.5, True)

# Right: eye-catching patient flow
rect(6.25, 1.98, 6.38, 4.96, panel, radius=True, border=line)
text(6.62, 2.22, 5.7, 0.3, "PATIENT PROFILE  →  EXPLAINABLE OUTPUT", 10, yellow, True)
text(6.62, 2.53, 5.5, 0.3, "A live run in four moves", 17, white, True)

# Patient input node
rect(6.65, 3.1, 1.48, 1.1, RGBColor(33, 57, 79), radius=True, border=blue)
text(6.78, 3.22, 1.22, 0.25, "01  CAPTURE", 9, blue, True, PP_ALIGN.CENTER)
text(6.78, 3.57, 1.22, 0.45, "Age • BP\nBMI • Lifestyle", 10, white, True, PP_ALIGN.CENTER)

# Eye / focus ring diagram
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.53), Inches(3.18), Inches(0.86), Inches(0.86))
circle.fill.solid(); circle.fill.fore_color.rgb = navy; circle.line.color.rgb = coral; circle.line.width = Pt(2.5)
inner = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.78), Inches(3.43), Inches(0.36), Inches(0.36))
inner.fill.solid(); inner.fill.fore_color.rgb = coral; inner.line.color.rgb = coral
text(8.26, 4.12, 1.42, 0.25, "02  VALIDATE", 9, coral, True, PP_ALIGN.CENTER)
line_shape(8.13, 3.65, 8.5, 3.65, blue, 2.2, True)

# Signal cards
cards = [
    (9.55, 3.03, "03  PREDICT", "27.4% CVD", cyan, "LOW"),
    (9.55, 4.05, "04  EXPLAIN", "Systolic BP", coral, "DRIVER"),
    (9.55, 5.07, "ACTION PLAN", "Personalised", yellow, "ACT"),
]
for x, y, label, value, color, tag in cards:
    rect(x, y, 2.68, 0.82, RGBColor(27, 49, 70), radius=True, border=line)
    rect(x, y, 0.08, 0.82, color, radius=True)
    text(x + 0.22, y + 0.1, 1.52, 0.18, label, 8, muted, True)
    text(x + 0.22, y + 0.32, 1.75, 0.28, value, 13, white, True)
    pill(x + 2.02, y + 0.28, 0.5, tag, color)

# Flow arrows and evidence callout
line_shape(8.14, 3.65, 9.45, 3.45, muted, 1.3, True)
line_shape(8.98, 4.05, 9.45, 4.45, muted, 1.3, True)
line_shape(8.98, 4.05, 9.45, 5.45, muted, 1.3, True)

rect(6.65, 5.92, 2.7, 0.58, RGBColor(30, 52, 71), radius=True, border=line)
text(6.82, 6.03, 2.35, 0.28, "Transparent by design", 10, white, True, PP_ALIGN.CENTER)
text(9.62, 6.18, 2.45, 0.28, "Inputs → models → evidence → action", 10, muted, False, PP_ALIGN.CENTER)

# Footer
line_shape(0.68, 7.09, 12.65, 7.09, line, 0.8)
text(0.7, 7.15, 8.8, 0.2, "Live demo: Streamlit dashboard • cached models • no patient data stored", 8, muted)
text(10.1, 7.15, 2.5, 0.2, "Research / education only", 8, yellow, True, PP_ALIGN.RIGHT)


def new_slide(section, title, subtitle, number):
    global slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = navy
    text(0.65, 0.28, 8.6, 0.28, f"{number:02d}  |  {section}", 10, yellow, True)
    text(0.65, 0.62, 11.8, 0.46, title, 23, white, True)
    text(0.67, 1.16, 11.6, 0.25, subtitle, 10, muted)
    line_shape(0.68, 7.09, 12.65, 7.09, line, 0.8)
    text(0.7, 7.15, 8.8, 0.2, "CARDIOLENS  |  Streamlit demo  |  research / education only", 8, muted)
    text(11.3, 7.15, 1.35, 0.2, f"{number:02d}", 8, yellow, True, PP_ALIGN.RIGHT)


# Slide 02: implementation map
new_slide(
    "WHAT IS IMPLEMENTED",
    "From captured vitals to an actionable care plan",
    "The app keeps the decision path visible: inputs, engineered features, model outputs, and recommendations.",
    2,
)
steps = [
    ("1", "CAPTURE", "Age, BP, height, weight, cholesterol, glucose, smoking, alcohol, activity."),
    ("2", "PREPARE", "BMI, pulse pressure, MAP, and combined metabolic flags are calculated."),
    ("3", "FORECAST", "XGBoost estimates overweight / obesity probability; LightGBM estimates CVD risk."),
    ("4", "EXPLAIN", "Feature contribution bars identify the strongest model drivers."),
    ("5", "ACT", "Risk-tiered wellness guidance becomes a personalized next-step plan."),
]
for i, (number, label, detail) in enumerate(steps):
    y = 1.78 + i * 0.87
    add_circle(0.82, y, 0.34, [blue, coral, cyan, violet, yellow][i], border=[blue, coral, cyan, violet, yellow][i])
    text(0.82, y + 0.01, 0.34, 0.3, number, 9, navy, True, PP_ALIGN.CENTER)
    text(1.35, y - 0.01, 1.1, 0.22, label, 10, [blue, coral, cyan, violet, yellow][i], True)
    text(2.55, y - 0.01, 3.45, 0.42, detail, 9, white)
    if i < len(steps) - 1:
        connector(0.99, y + 0.35, 0.99, y + 0.84, line, 1.2, True)

rect(6.45, 1.78, 5.85, 4.75, panel, radius=True, border=line)
text(6.8, 2.05, 5.1, 0.25, "THE DEMO PROMISE", 10, cyan, True)
text(6.8, 2.38, 4.9, 0.42, "Every score has a visible reason\nand a practical next step.", 18, white, True)
metrics = [("02", "MODEL ARTIFACTS", "cached and loaded at runtime", coral),
           ("13", "STAGE 2 FEATURES", "clinical + engineered signals", cyan),
           ("01", "ACTION PLAN", "tailored to the risk tier", yellow)]
for i, (value, label, detail, accent) in enumerate(metrics):
    y = 3.25 + i * 0.88
    text(6.82, y, 0.65, 0.34, value, 21, accent, True)
    text(7.65, y + 0.01, 2.1, 0.2, label, 9, white, True)
    text(7.65, y + 0.26, 3.65, 0.2, detail, 9, muted)

# Slide 03: interactive flow
new_slide(
    "EYE-CATCHING FLOW",
    "The interactive capture-to-forecast loop",
    "A presenter can demonstrate the full loop quickly with the live sidebar controls and instant recalculation.",
    3,
)
flow = [(0.95, "01", "PROFILE", "Set age, BP, BMI\nand lifestyle", blue),
        (3.35, "02", "ENGINEER", "Calculate BMI, MAP\nand pulse pressure", coral),
        (5.75, "03", "CASCADE", "Pass metabolic risk\ninto CVD model", cyan),
        (8.15, "04", "INSPECT", "Read risk drivers\nand feature weights", violet),
        (10.55, "05", "ACT", "Open the tiered\nwellness plan", yellow)]
for i, (x, number, label, detail, accent) in enumerate(flow):
    add_circle(x, 2.18, 0.78, panel_2, border=accent, width=2.2)
    text(x, 2.37, 0.78, 0.32, number, 13, accent, True, PP_ALIGN.CENTER)
    text(x - 0.25, 3.12, 1.28, 0.22, label, 10, accent, True, PP_ALIGN.CENTER)
    text(x - 0.45, 3.46, 1.68, 0.48, detail, 9, white, False, PP_ALIGN.CENTER)
    if i < len(flow) - 1:
        connector(x + 0.84, 2.57, x + 2.28, 2.57, muted, 1.4, True)
rect(0.88, 4.55, 11.65, 1.42, panel, radius=True, border=line)
text(1.2, 4.82, 2.1, 0.22, "PRESENTER CUE", 9, yellow, True)
text(1.2, 5.13, 10.7, 0.45, "Start with a high-pressure profile, show the risk tier, open the feature contribution chart, then change one lifestyle signal and recalculate.", 15, white, True)
text(1.2, 5.64, 10.7, 0.2, "The audience sees the same patient profile move from input to evidence to action.", 9, muted)

# Slide 04: 3D visual explanation
new_slide(
    "3D VISUAL EXPLANATION",
    "A clinical signal map for the three-layer engine",
    "The 3D stack is a visual metaphor for how raw biometrics become a forecast and then an interpretable plan.",
    4,
)
layer_data = [(1.0, 2.0, 5.15, coral, "LAYER 01", "BIOMETRICS", "patient inputs captured in the sidebar"),
             (1.62, 3.28, 5.15, cyan, "LAYER 02", "FEATURE SPACE", "BMI, MAP, pulse pressure, metabolic flags"),
             (2.24, 4.56, 5.15, yellow, "LAYER 03", "FORECAST SURFACE", "CVD probability, drivers, and action plan")]
for x, y, width, accent, kicker, label, detail in layer_data:
    rect(x + 0.18, y + 0.18, width, 0.92, RGBColor(5, 13, 24), radius=True)
    rect(x + 0.08, y + 0.08, width, 0.92, accent, radius=True)
    rect(x, y, width, 0.92, panel, radius=True, border=accent)
    text(x + 0.3, y + 0.14, 1.0, 0.18, kicker, 8, accent, True)
    text(x + 1.35, y + 0.11, 2.0, 0.24, label, 12, white, True)
    text(x + 1.35, y + 0.46, 3.2, 0.2, detail, 9, muted)
connector(3.58, 3.0, 3.58, 3.23, muted, 1.5, True)
connector(4.2, 4.28, 4.2, 4.51, muted, 1.5, True)
rect(8.2, 2.08, 3.9, 3.65, panel, radius=True, border=line)
text(8.58, 2.36, 3.1, 0.22, "THE MODEL LENS", 10, violet, True)
add_circle(9.45, 2.95, 1.72, navy, coral, 2.7)
add_circle(9.78, 3.28, 1.06, navy, cyan, 2.0)
add_circle(10.14, 3.64, 0.34, coral, coral, 1)
text(8.62, 4.95, 3.05, 0.28, "Risk is inspectable", 15, white, True, PP_ALIGN.CENTER)
text(8.62, 5.3, 3.05, 0.22, "Feature drivers turn a probability into a conversation.", 9, muted, False, PP_ALIGN.CENTER)

# Slide 05: contrasting demo stories
new_slide(
    "VERIFIED DEMO",
    "One app, two contrasting patient stories",
    "Use these illustrative profiles to show how the same interface responds to different cardiometabolic signals.",
    5,
)
profiles = [(0.9, "OPTIMAL BASELINE", "45-year-old profile", cyan, "BMI 22.4", "BP 115 / 75", "Active  |  non-smoker", "LOWER SIGNAL LOAD"),
            (6.82, "ELEVATED PROFILE", "58-year-old profile", coral, "BMI 33.0", "BP 150 / 95", "Smoker  |  sedentary", "HIGHER SIGNAL LOAD")]
for x, heading, age_text, accent, bmi_text, bp_text, lifestyle, status in profiles:
    rect(x, 1.85, 5.55, 4.55, panel, radius=True, border=accent)
    text(x + 0.35, 2.14, 4.8, 0.22, heading, 10, accent, True)
    text(x + 0.35, 2.48, 4.8, 0.3, age_text, 16, white, True)
    text(x + 0.35, 2.98, 4.75, 0.22, "INPUT SIGNALS", 8, muted, True)
    for j, value in enumerate([bmi_text, bp_text, lifestyle]):
        y = 3.35 + j * 0.42
        rect(x + 0.35, y, 4.75, 0.3, panel_2, radius=True, border=line)
        text(x + 0.55, y + 0.02, 4.3, 0.24, value, 10, white, True)
    text(x + 0.35, 4.78, 2.15, 0.18, "STAGE 2 CVD RISK", 8, muted, True)
    text(x + 0.35, 5.05, 2.0, 0.42, "27.4%" if accent == cyan else "HIGHER", 22, accent, True)
    pill(x + 3.58, 5.08, 1.18, status, accent)
    text(x + 0.35, 5.72, 4.65, 0.22, "Then inspect drivers and open the action plan.", 9, muted)

# Slide 06: explainability and advice
new_slide(
    "EXPLAINABILITY + ACTION",
    "The result is not just a score",
    "CardioLens translates the forecast into visible drivers, a risk tier, and a tailored wellness response.",
    6,
)
cards = [(0.92, "SCORE", "27.4%", "10-year CVD estimate", cyan),
         (3.55, "EXPLAIN", "BP + BMI", "leading model signals", coral),
         (6.18, "TIER", "LOW", "green status indicator", yellow),
         (8.81, "ACT", "PLAN", "prioritized next steps", violet)]
for x, label, value, detail, accent in cards:
    rect(x, 1.95, 2.25, 1.42, panel, radius=True, border=line)
    rect(x, 1.95, 0.08, 1.42, accent, radius=True)
    text(x + 0.25, 2.15, 1.65, 0.18, label, 8, muted, True)
    text(x + 0.25, 2.45, 1.75, 0.35, value, 18, accent, True)
    text(x + 0.25, 2.93, 1.7, 0.18, detail, 8, white)
connector(1.98, 3.73, 10.0, 3.73, line, 1.3, True)
rect(0.92, 4.12, 11.05, 1.46, panel, radius=True, border=line)
text(1.25, 4.4, 2.05, 0.22, "ACTION PLAN CONTRACT", 9, yellow, True)
text(1.25, 4.78, 10.1, 0.3, "Protective factors  •  factors requiring attention  •  prioritized lifestyle guidance", 14, white, True)
text(1.25, 5.2, 10.1, 0.2, "The recommendation engine adapts to BP, BMI, smoking, alcohol, activity, cholesterol, glucose, and risk tier.", 9, muted)

# Slide 07: close
new_slide(
    "DEMO CLOSE",
    "A clear path from signal to next step",
    "CardioLens makes a multi-stage machine-learning forecast understandable, inspectable, and useful in a live walkthrough.",
    7,
)
close_steps = [("01", "Start with the profile", "Capture vitals and lifestyle inputs in the sidebar.", blue),
               ("02", "Trace the cascade", "Show how metabolic probability feeds the cardiovascular model.", cyan),
               ("03", "End with action", "Open the feature drivers and personalized wellness plan.", yellow)]
for i, (number, heading, detail, accent) in enumerate(close_steps):
    y = 1.95 + i * 1.25
    add_circle(0.95, y, 0.62, accent, accent)
    text(0.95, y + 0.13, 0.62, 0.32, number, 11, navy, True, PP_ALIGN.CENTER)
    text(1.92, y + 0.02, 3.6, 0.25, heading, 15, white, True)
    text(1.92, y + 0.36, 4.5, 0.25, detail, 10, muted)
    if i < len(close_steps) - 1:
        connector(1.26, y + 0.65, 1.26, y + 1.18, line, 1.3, True)
rect(7.55, 2.0, 4.35, 3.45, panel, radius=True, border=cyan)
text(7.95, 2.35, 3.55, 0.22, "CARDIOLENS IN ONE LINE", 9, cyan, True)
text(7.95, 2.84, 3.45, 1.15, "Inputs\n→ models\n→ evidence\n→ action", 24, white, True, PP_ALIGN.CENTER)
text(7.95, 4.55, 3.45, 0.35, "For research and educational use only.", 10, yellow, True, PP_ALIGN.CENTER)

prs.save(OUT)
print(OUT)
